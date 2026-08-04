# -*- coding: utf-8 -*-
"""
Genera la presentacion PPTX de InvenPro:
BD, frontend, backend, tecnologias, SOLID y POO.
Uso:  py gen_presentacion_pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- Paleta ----
INK    = RGBColor(0x0F, 0x17, 0x2A)
ACCENT = RGBColor(0x25, 0x63, 0xEB)
SKY    = RGBColor(0x0E, 0xA5, 0xE9)
LIGHT  = RGBColor(0xF1, 0xF5, 0xF9)
CARD   = RGBColor(0xFF, 0xFF, 0xFF)
MUTED  = RGBColor(0x64, 0x74, 0x8B)
GOOD   = RGBColor(0x16, 0xA3, 0x4A)
AMBER  = RGBColor(0xD9, 0x77, 0x06)
RED    = RGBColor(0xDC, 0x26, 0x26)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height

def slide():
    return prs.slides.add_slide(BLANK)

def rect(s, l, t, w, h, fill, line=None, shape=MSO_SHAPE.RECTANGLE, radius=None):
    sp = s.shapes.add_shape(shape, l, t, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp

def text(s, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    """runs: lista de parrafos; cada parrafo es lista de (txt, size, color, bold)."""
    tb = s.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(4); tf.margin_top = tf.margin_bottom = Pt(2)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        for (txt, size, color, bold) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold
            r.font.name = "Calibri"
    return tb

def header(s, kicker, title, color=ACCENT):
    rect(s, 0, 0, SW, Inches(1.15), INK)
    rect(s, 0, Inches(1.15), SW, Pt(4), color)
    text(s, Inches(0.6), Inches(0.16), Inches(12), Inches(0.35),
         [[(kicker.upper(), 12, SKY, True)]])
    text(s, Inches(0.6), Inches(0.42), Inches(12), Inches(0.62),
         [[(title, 26, WHITE, True)]])

def bullets(s, l, t, w, items, size=15, gap=True):
    runs = []
    for it in items:
        if isinstance(it, tuple):
            head, sub = it
            runs.append([("•  ", size, ACCENT, True), (head+"  ", size, INK, True), (sub, size-1, MUTED, False)])
        else:
            runs.append([("•  ", size, ACCENT, True), (it, size, INK, False)])
    tb = text(s, l, t, w, Inches(5), runs)
    for p in tb.text_frame.paragraphs:
        p.space_after = Pt(10 if gap else 4); p.line_spacing = 1.05
    return tb

def card(s, l, t, w, h, title, lines, accent=ACCENT, title_color=None):
    rect(s, l, t, w, h, CARD, line=RGBColor(0xE2,0xE8,0xF0))
    rect(s, l, t, Inches(0.09), h, accent)
    text(s, l+Inches(0.25), t+Inches(0.14), w-Inches(0.4), Inches(0.4),
         [[(title, 15, title_color or INK, True)]])
    runs = [[(ln, 12.5, MUTED, False)] for ln in lines]
    tb = text(s, l+Inches(0.25), t+Inches(0.62), w-Inches(0.42), h-Inches(0.7), runs)
    for p in tb.text_frame.paragraphs:
        p.space_after = Pt(5); p.line_spacing = 1.02
    return tb

def footer(s, n):
    text(s, Inches(0.6), SH-Inches(0.42), Inches(9), Inches(0.3),
         [[("InvenPro · Sistema de Inventario y POS", 9, MUTED, False)]])
    text(s, SW-Inches(1.2), SH-Inches(0.42), Inches(0.7), Inches(0.3),
         [[(str(n), 9, MUTED, False)]], align=PP_ALIGN.RIGHT)

# =================================================================
# 1) PORTADA
# =================================================================
s = slide()
rect(s, 0, 0, SW, SH, INK)
rect(s, 0, Inches(4.55), SW, Pt(5), ACCENT)
rect(s, Inches(0.9), Inches(1.5), Inches(1.5), Inches(1.5), ACCENT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, Inches(0.9), Inches(1.5), Inches(1.5), Inches(1.5),
     [[("IP", 44, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, Inches(0.9), Inches(3.25), Inches(11.5), Inches(1.1),
     [[("InvenPro", 54, WHITE, True)]])
text(s, Inches(0.92), Inches(4.65), Inches(11.5), Inches(0.6),
     [[("Sistema de Inventario y Punto de Venta", 22, SKY, False)]])
text(s, Inches(0.92), Inches(5.35), Inches(11.5), Inches(0.8),
     [[("Arquitectura, tecnologias y diseno orientado a objetos", 16, RGBColor(0xCB,0xD5,0xE1), False)]])
text(s, Inches(0.92), Inches(6.6), Inches(11.5), Inches(0.4),
     [[("Documento tecnico del proyecto", 12, MUTED, False)]])

# =================================================================
# 2) AGENDA
# =================================================================
s = slide(); header(s, "Contenido", "Agenda")
items = [
    ("Arquitectura general", "vision de las tres capas"),
    ("Base de datos", "PostgreSQL sobre Supabase"),
    ("Frontend", "React 18 + Tailwind, sin bundler"),
    ("Backend", "Supabase: Auth, Realtime, Edge Functions"),
    ("Stack tecnologico", "librerias y su proposito"),
    ("Programacion Orientada a Objetos", "modelos y servicios"),
    ("Principios SOLID", "donde se aplicaron"),
    ("Patrones de diseno", "Repository, Observer, Adapter"),
]
left = [it for it in items[:4]]; right = [it for it in items[4:]]
bullets(s, Inches(0.7), Inches(1.6), Inches(6), left, size=17)
bullets(s, Inches(6.9), Inches(1.6), Inches(6), right, size=17)
footer(s, 2)

# =================================================================
# 3) ARQUITECTURA GENERAL
# =================================================================
s = slide(); header(s, "Vision global", "Arquitectura general")
y = Inches(2.0); bh = Inches(2.2); bw = Inches(3.7)
card(s, Inches(0.55), y, bw, bh, "Cliente (Navegador)",
     ["React 18 + Babel Standalone", "Tailwind CSS · Chart.js", "SheetJS · html5-qrcode",
      "Componentes por dominio", "JSX compilado en el navegador"], accent=ACCENT)
card(s, Inches(4.8), y, bw, bh, "Supabase (Backend)",
     ["API REST (PostgREST)", "Auth · Realtime", "Edge Functions (Deno)", "RPC · pg_cron · Vault",
      "Capa de servicios (DB.*)"], accent=SKY)
card(s, Inches(9.05), y, bw, bh, "PostgreSQL (Datos)",
     ["14 tablas + 4 vistas", "Row Level Security", "Funciones PL/pgSQL", "Realtime (postgres_changes)",
      "Indices y restricciones"], accent=GOOD)
# flechas
for cx in [Inches(4.42), Inches(8.67)]:
    a = rect(s, cx, y+Inches(0.9), Inches(0.42), Inches(0.42), INK, shape=MSO_SHAPE.RIGHT_ARROW)
text(s, Inches(0.55), Inches(4.5), Inches(12.2), Inches(0.6),
     [[("Flujo: ", 14, INK, True),
       ("el navegador consume la capa de servicios de Supabase; los cambios se propagan en vivo por Realtime a todos los clientes.", 14, MUTED, False)]])
text(s, Inches(0.55), Inches(5.25), Inches(12.2), Inches(1.4),
     [[("Sin servidor propio ni bundler: ", 13.5, INK, True),
       ("la logica de negocio vive en clases del lado cliente (data.js) y en funciones/RPC del lado servidor. Es una arquitectura BaaS (Backend as a Service).", 13.5, MUTED, False)]])
footer(s, 3)

# =================================================================
# 4) BASE DE DATOS
# =================================================================
s = slide(); header(s, "Persistencia", "Base de datos")
card(s, Inches(0.55), Inches(1.55), Inches(3.9), Inches(2.55), "Motor",
     ["PostgreSQL", "Gestionado por Supabase (BaaS)", "API REST autogenerada (PostgREST)",
      "Acceso via cliente supabase-js"], accent=GOOD)
card(s, Inches(4.7), Inches(1.55), Inches(3.9), Inches(2.55), "Estructura",
     ["10 tablas de negocio", "productos, facturas, ingresos,", "turnos, proveedores, cajeros...",
      "4 vistas de agregacion", "ventas_mes, top_productos..."], accent=ACCENT)
card(s, Inches(8.85), Inches(1.55), Inches(3.9), Inches(2.55), "Caracteristicas",
     ["Row Level Security (RLS)", "RPC: increment_stock /", "decrement_stock (atomico)",
      "Realtime (postgres_changes)", "Hash SHA-256 de contrasenas"], accent=SKY)
text(s, Inches(0.55), Inches(4.5), Inches(12.2), Inches(1.6),
     [[("Integridad y concurrencia:  ", 14, INK, True),
       ("el ajuste de inventario se hace con funciones RPC en la base (no en el cliente), garantizando operaciones atomicas al vender o ingresar mercancia. Un indice unico parcial asegura un solo turno abierto por cajero.", 14, MUTED, False)]])
footer(s, 4)

# =================================================================
# 5) FRONTEND
# =================================================================
s = slide(); header(s, "Capa de presentacion", "Frontend")
bullets(s, Inches(0.7), Inches(1.55), Inches(6.1), [
    ("React 18 (via CDN):", "UI declarativa por componentes"),
    ("Babel Standalone:", "compila JSX en el navegador, sin bundler"),
    ("Tailwind CSS v3 (prefijo tw-):", "diseno responsive"),
    ("CSS propio:", "tokens y tema claro / oscuro"),
    ("Hooks:", "useState, useMemo, useEffect, useRealtimeSync"),
], size=15)
card(s, Inches(7.1), Inches(1.55), Inches(5.65), Inches(2.15), "Librerias de apoyo",
     ["Chart.js 4  —  graficos del dashboard y reportes",
      "SheetJS (xlsx)  —  exportar a Excel",
      "html5-qrcode  —  escaner de codigos de barras"], accent=SKY)
card(s, Inches(7.1), Inches(3.85), Inches(5.65), Inches(2.4), "Componentes por dominio",
     ["Dashboard · Inventario · Ingreso",
      "Vencimientos · Proveedores · Cajeros",
      "POS (cajero) · Reportes · Ajustes",
      "Cada modulo en su propio archivo .jsx"], accent=ACCENT)
text(s, Inches(0.7), Inches(4.7), Inches(6.1), Inches(1.5),
     [[("Estrategia responsive: ", 13.5, INK, True),
       ("utilidades Tailwind en el JSX (colapsar columnas, mostrar/ocultar, tarjetas en movil) sobre una base de CSS para estructura y tokens.", 13.5, MUTED, False)]])
footer(s, 5)

# =================================================================
# 6) BACKEND
# =================================================================
s = slide(); header(s, "Capa de servicios", "Backend  —  Supabase")
card(s, Inches(0.55), Inches(1.55), Inches(3.9), Inches(2.5), "Datos y API",
     ["PostgreSQL + PostgREST", "API REST autogenerada", "RPC para logica transaccional",
      "(increment/decrement stock)"], accent=GOOD)
card(s, Inches(4.7), Inches(1.55), Inches(3.9), Inches(2.5), "Auth y Realtime",
     ["Login contra usuarios_sistema", "Contrasenas con hash SHA-256", "Realtime: canal unico",
      "postgres_changes en las tablas"], accent=ACCENT)
card(s, Inches(8.85), Inches(1.55), Inches(3.9), Inches(2.5), "Edge Functions (Deno)",
     ["enviar-alerta (TypeScript)", "Correo de vencimientos", "Gmail SMTP (denomailer)",
      "pg_cron + pg_net (envio diario)", "Vault para secretos"], accent=SKY)
text(s, Inches(0.55), Inches(4.45), Inches(12.2), Inches(1.6),
     [[("Backend as a Service: ", 14, INK, True),
       ("no hay servidor propio que mantener. La logica se reparte entre la capa de servicios del cliente (clases en data.js), las funciones de la base (RPC) y las Edge Functions serverless para tareas fuera del navegador (envio de correos, cron).", 14, MUTED, False)]])
footer(s, 6)

# =================================================================
# 7) STACK TECNOLOGICO (tabla)
# =================================================================
s = slide(); header(s, "Resumen", "Stack tecnologico")
rows = [
    ("Capa", "Tecnologia", "Proposito"),
    ("Base de datos", "PostgreSQL (Supabase)", "Persistencia, vistas, RPC, RLS"),
    ("Backend", "Supabase (Auth, Realtime, Storage)", "API REST, autenticacion, sincronizacion"),
    ("Serverless", "Edge Functions (Deno / TypeScript)", "Envio de correos de alerta"),
    ("Automatizacion", "pg_cron + pg_net + Vault", "Alertas diarias programadas"),
    ("Frontend", "React 18 + Babel Standalone", "UI por componentes, sin bundler"),
    ("Estilos", "Tailwind CSS v3 + CSS propio", "Responsive y tema claro/oscuro"),
    ("Librerias", "Chart.js · SheetJS · html5-qrcode", "Graficos, Excel, escaner"),
    ("Lenguajes", "JavaScript (ES6+), SQL, TypeScript", "Cliente, base de datos, serverless"),
]
tw = Inches(12.2); th = Inches(5.0)
tbl = s.shapes.add_table(len(rows), 3, Inches(0.55), Inches(1.5), tw, th).table
tbl.columns[0].width = Inches(2.5); tbl.columns[1].width = Inches(4.5); tbl.columns[2].width = Inches(5.2)
for ci, val in enumerate(rows[0]):
    c = tbl.cell(0, ci); c.fill.solid(); c.fill.fore_color.rgb = INK
    tf = c.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; p.text = val
    p.runs[0].font.bold = True; p.runs[0].font.size = Pt(14); p.runs[0].font.color.rgb = WHITE
for ri in range(1, len(rows)):
    for ci, val in enumerate(rows[ri]):
        c = tbl.cell(ri, ci); c.fill.solid()
        c.fill.fore_color.rgb = CARD if ri % 2 else LIGHT
        tf = c.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; p.text = val
        p.runs[0].font.size = Pt(12.5)
        p.runs[0].font.bold = (ci == 0)
        p.runs[0].font.color.rgb = INK if ci == 0 else MUTED
footer(s, 7)

# =================================================================
# 8) POO
# =================================================================
s = slide(); header(s, "Diseno", "Programacion Orientada a Objetos")
text(s, Inches(0.6), Inches(1.4), Inches(12.2), Inches(0.5),
     [[("La capa de dominio (data.js) encapsula entidades y acceso a datos con clases ES6.", 15, INK, True)]])
card(s, Inches(0.55), Inches(2.05), Inches(4.0), Inches(3.0), "6 Modelos (entidades)",
     ["Producto · Usuario · Cajero", "Proveedor · Turno · Factura", "",
      "Encapsulan datos + logica de", "dominio con getters calculados:",
      "stockBajo, valorTotal, esAdmin,", "abierto, initials, cantidadItems"], accent=ACCENT)
card(s, Inches(4.7), Inches(2.05), Inches(4.0), Inches(3.0), "8 Servicios",
     ["AuthService · ProductoService", "FacturaService · TurnoService", "CajeroService · ProveedorService",
      "IngresoService · ConfigService", "", "Encapsulan las operaciones", "de base de datos (metodos async)"], accent=SKY)
card(s, Inches(8.85), Inches(2.05), Inches(3.9), Inches(3.0), "Infraestructura",
     ["DataStore  —  agrega e hidrata", "el estado de la aplicacion", "",
      "EventBus  —  eventos (pub/sub)", "RealtimeManager  —  sincroniza", "",
      "camelize / snakify  —  adaptan", "snake_case (BD) a camelCase"], accent=GOOD)
text(s, Inches(0.55), Inches(5.3), Inches(12.2), Inches(1.4),
     [[("Conceptos aplicados:  ", 13.5, INK, True),
       ("encapsulacion (constructor + this), abstraccion (los metodos ocultan las llamadas a Supabase), propiedades calculadas (get) e instanciacion uniforme (camelize -> new Modelo).", 13.5, MUTED, False)]])
footer(s, 8)

# =================================================================
# 9) SOLID (tabla)
# =================================================================
s = slide(); header(s, "Buenas practicas", "Principios SOLID  —  donde se aplicaron")
rows = [
    ("Principio", "Aplicacion en InvenPro"),
    ("S — Responsabilidad unica",
     "Cada servicio tiene una sola razon de cambio (AuthService=login, ProductoService=productos...). "
     "Los modelos solo modelan datos+dominio. El panel admin se dividio en un archivo por modulo."),
    ("O — Abierto / Cerrado",
     "Se agregan nuevos modelos o servicios sin modificar los existentes. Los presets de IA y la deteccion "
     "de proveedor se extienden por configuracion, no tocando el codigo."),
    ("L — Sustitucion de Liskov",
     "camelize produce instancias uniformes; cualquier Producto o Factura es intercambiable donde se espera "
     "esa entidad, sin romper el comportamiento."),
    ("I — Segregacion de interfaces",
     "API por espacios de nombres (DB.auth, DB.productos, DB.turnos...). Cada componente usa solo el servicio "
     "que necesita, no un objeto monolitico."),
    ("D — Inversion de dependencias",
     "Los componentes dependen de la abstraccion window.DB (capa de servicios) y de EventBus, no de llamadas "
     "directas a Supabase."),
]
tbl = s.shapes.add_table(len(rows), 2, Inches(0.55), Inches(1.5), Inches(12.2), Inches(5.2)).table
tbl.columns[0].width = Inches(3.3); tbl.columns[1].width = Inches(8.9)
colors = [ACCENT, GOOD, AMBER, SKY, RED]
for ci, val in enumerate(rows[0]):
    c = tbl.cell(0, ci); c.fill.solid(); c.fill.fore_color.rgb = INK
    p = c.text_frame.paragraphs[0]; p.text = val
    p.runs[0].font.bold = True; p.runs[0].font.size = Pt(14); p.runs[0].font.color.rgb = WHITE
for ri in range(1, len(rows)):
    c0 = tbl.cell(ri, 0); c0.fill.solid(); c0.fill.fore_color.rgb = LIGHT
    p = c0.text_frame.paragraphs[0]; p.text = rows[ri][0]
    p.runs[0].font.bold = True; p.runs[0].font.size = Pt(12.5); p.runs[0].font.color.rgb = colors[ri-1]
    c1 = tbl.cell(ri, 1); c1.fill.solid(); c1.fill.fore_color.rgb = CARD
    c1.text_frame.word_wrap = True
    p = c1.text_frame.paragraphs[0]; p.text = rows[ri][1]
    p.runs[0].font.size = Pt(11.5); p.runs[0].font.color.rgb = MUTED
footer(s, 9)

# =================================================================
# 10) PATRONES
# =================================================================
s = slide(); header(s, "Arquitectura", "Patrones de diseno")
card(s, Inches(0.55), Inches(1.6), Inches(6.0), Inches(1.55), "Repository / Facade",
     ["DataStore + capa de servicios (DB) abstraen el acceso",
      "a datos: los componentes no consultan la BD directamente."], accent=ACCENT)
card(s, Inches(6.75), Inches(1.6), Inches(6.0), Inches(1.55), "Observer / Pub-Sub",
     ["EventBus (emit / on) desacopla la sincronizacion Realtime",
      "de los componentes que reaccionan a los cambios."], accent=GOOD)
card(s, Inches(0.55), Inches(3.35), Inches(6.0), Inches(1.55), "Adapter",
     ["camelize / snakify traducen entre snake_case (base de",
      "datos) y camelCase (JavaScript) de forma transparente."], accent=AMBER)
card(s, Inches(6.75), Inches(3.35), Inches(6.0), Inches(1.55), "Singleton / Service Layer",
     ["Cliente Supabase unico (window.db) y RealtimeManager con",
      "un solo canal; un servicio por entidad del dominio."], accent=SKY)
text(s, Inches(0.55), Inches(5.25), Inches(12.2), Inches(1.2),
     [[("Resultado: ", 14, INK, True),
       ("codigo modular, con responsabilidades separadas y bajo acoplamiento, facil de mantener y extender pese a no usar bundler ni framework de build.", 14, MUTED, False)]])
footer(s, 10)

# =================================================================
# 11) CIERRE
# =================================================================
s = slide()
rect(s, 0, 0, SW, SH, INK)
rect(s, 0, Inches(3.05), SW, Pt(5), ACCENT)
text(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.0),
     [[("Resumen", 40, WHITE, True)]])
text(s, Inches(0.92), Inches(3.35), Inches(11.6), Inches(2.6), [
    [("Base de datos:  ", 17, SKY, True), ("PostgreSQL sobre Supabase (vistas, RPC, RLS, Realtime).", 17, RGBColor(0xCB,0xD5,0xE1), False)],
    [("Frontend:  ", 17, SKY, True), ("React 18 + Babel + Tailwind, sin bundler.", 17, RGBColor(0xCB,0xD5,0xE1), False)],
    [("Backend:  ", 17, SKY, True), ("Supabase (Auth, Realtime, Edge Functions en Deno).", 17, RGBColor(0xCB,0xD5,0xE1), False)],
    [("POO:  ", 17, SKY, True), ("6 modelos + 8 servicios + DataStore en data.js.", 17, RGBColor(0xCB,0xD5,0xE1), False)],
    [("SOLID:  ", 17, SKY, True), ("responsabilidad unica, interfaces segregadas e inversion de dependencias.", 17, RGBColor(0xCB,0xD5,0xE1), False)],
])
for p in s.shapes[-1].text_frame.paragraphs:
    p.space_after = Pt(12)
text(s, Inches(0.92), Inches(6.7), Inches(11.5), Inches(0.5),
     [[("InvenPro · Documento tecnico del proyecto", 12, MUTED, False)]])

out = "InvenPro_Presentacion.pptx"
prs.save(out)
print("OK ->", out, "|", len(prs.slides._sldIdLst), "diapositivas")
