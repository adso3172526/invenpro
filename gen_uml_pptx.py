"""
Genera diagrama de clases UML con RELACIONES para InvenPro.
Cada clase muestra atributos + métodos, conectadas con flechas.
Ejecutar: py gen_uml_pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── Colores ──
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BG          = RGBColor(0xF8, 0xF9, 0xFB)
DARK        = RGBColor(0x1E, 0x29, 0x3B)
GRAY        = RGBColor(0x64, 0x74, 0x8B)
LGRAY       = RGBColor(0xE2, 0xE8, 0xF0)
ACCENT      = RGBColor(0x37, 0x5B, 0xE9)
ACCENT_SOFT = RGBColor(0xEB, 0xEF, 0xFD)
GREEN       = RGBColor(0x16, 0xA3, 0x4A)
GREEN_SOFT  = RGBColor(0xDC, 0xFC, 0xE7)
PURPLE      = RGBColor(0x7C, 0x3A, 0xED)
PURPLE_SOFT = RGBColor(0xED, 0xE9, 0xFE)
ORANGE      = RGBColor(0xEA, 0x58, 0x0C)
ORANGE_SOFT = RGBColor(0xFF, 0xED, 0xD5)
TEAL        = RGBColor(0x0D, 0x94, 0x88)
TEAL_SOFT   = RGBColor(0xCC, 0xFB, 0xF1)
BORDER      = RGBColor(0xCB, 0xD5, 0xE1)
STRIPE      = RGBColor(0xF1, 0xF5, 0xF9)

FONT = "Geist"
MONO = "Geist Mono"

prs = Presentation()
prs.slide_width  = Inches(16)
prs.slide_height = Inches(9)

# ═══════════════════════════════════════
#  HELPERS
# ═══════════════════════════════════════

def rect(sl, l, t, w, h, fill=WHITE, brd=None, r=0.03):
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill; s.shadow.inherit = False
    if brd: s.line.color.rgb = brd; s.line.width = Pt(1)
    else: s.line.fill.background()
    if len(s.adjustments) > 0: s.adjustments[0] = r
    return s

def txt(sl, l, t, w, h, text, sz=10, bold=False, clr=DARK, fn=FONT, al=PP_ALIGN.LEFT):
    tb = sl.shapes.add_textbox(l, t, w, h); tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(sz)
    p.font.bold = bold; p.font.color.rgb = clr; p.font.name = fn; p.alignment = al
    return tb

def mtxt(sl, l, t, w, lines, sz=8, clr=DARK, fn=MONO):
    h = max(len(lines) * Pt(sz + 5), Pt(14))
    tb = sl.shapes.add_textbox(l, t, w, h); tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(sz); p.font.color.rgb = clr
        p.font.name = fn; p.space_after = Pt(1); p.space_before = Pt(0)
    return tb

def line(sl, x1, y1, x2, y2, clr=BORDER, w=1.0):
    c = sl.shapes.add_connector(1, x1, y1, x2, y2)
    c.line.color.rgb = clr; c.line.width = Pt(w)
    return c

def arrow(sl, x1, y1, x2, y2, clr=GRAY, w=1.2, dashed=False, head='triangle', tail=None):
    """Flecha con punta configurable."""
    c = sl.shapes.add_connector(1, x1, y1, x2, y2)
    c.line.color.rgb = clr; c.line.width = Pt(w)
    if dashed: c.line.dash_style = 4
    ln = c.line._ln
    if head:
        el = ln.makeelement(qn('a:tailEnd'), {'type': head, 'w': 'med', 'len': 'med'})
        ln.append(el)
    if tail:
        el = ln.makeelement(qn('a:headEnd'), {'type': tail, 'w': 'med', 'len': 'med'})
        ln.append(el)
    return c

def diamond_arrow(sl, x1, y1, x2, y2, clr=GRAY, w=1.5):
    """Flecha de composición ◆──→."""
    return arrow(sl, x1, y1, x2, y2, clr=clr, w=w, tail='diamond')

def cbox(sl, l, t, w, name, stereo, attrs, meths, hf, ht=WHITE, brd=BORDER, fs=8):
    """Caja UML de 3 secciones. Retorna (x_center, y_top, y_bottom, width, height)."""
    LH = Pt(fs + 4.5); PAD = Pt(5)
    HH = Pt(34) if stereo else Pt(26)
    ah = max(len(attrs), 1) * LH + PAD * 2
    mh = max(len(meths), 1) * LH + PAD * 2
    th = HH + ah + mh
    rect(sl, l, t, w, th, fill=WHITE, brd=brd)
    rect(sl, l, t, w, HH, fill=hf, brd=brd)
    if stereo:
        txt(sl, l, t + Pt(2), w, Pt(11), f"\u00AB{stereo}\u00BB", sz=6.5, clr=ht, al=PP_ALIGN.CENTER, fn=MONO)
        txt(sl, l, t + Pt(13), w, Pt(17), name, sz=10.5, bold=True, clr=ht, al=PP_ALIGN.CENTER)
    else:
        txt(sl, l, t + Pt(4), w, Pt(18), name, sz=10.5, bold=True, clr=ht, al=PP_ALIGN.CENTER)
    ya = t + HH
    if attrs: mtxt(sl, l + Pt(6), ya + PAD, w - Pt(12), attrs, sz=fs, clr=DARK)
    else: txt(sl, l + Pt(6), ya + PAD, w - Pt(12), LH, "(stateless)", sz=fs, clr=LGRAY, fn=MONO)
    ys = ya + ah
    line(sl, l + Pt(4), ys, l + w - Pt(4), ys, clr=LGRAY, w=0.5)
    if meths: mtxt(sl, l + Pt(6), ys + PAD, w - Pt(12), meths, sz=fs, clr=GRAY)
    else: txt(sl, l + Pt(6), ys + PAD, w - Pt(12), LH, "(sin metodos)", sz=fs, clr=LGRAY, fn=MONO)
    cx = l + w // 2
    return (l, t, l + w, t + th, cx, th)

def stitle(sl, title, sub=""):
    txt(sl, Inches(0.5), Inches(0.22), Inches(13), Inches(0.42), title, sz=21, bold=True, clr=DARK)
    if sub: txt(sl, Inches(0.5), Inches(0.66), Inches(14), Inches(0.28), sub, sz=10.5, clr=GRAY)

def nslide(bg=BG):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid(); sl.background.fill.fore_color.rgb = bg
    return sl

def label(sl, x, y, text, sz=7, clr=GRAY):
    """Etiqueta de relación junto a una flecha."""
    txt(sl, x, y, Inches(1.5), Pt(12), text, sz=sz, clr=clr, fn=MONO)


# ═══════════════════════════════════════════════════
#  SLIDE 1 — Portada
# ═══════════════════════════════════════════════════
s = nslide(DARK)
txt(s, Inches(2), Inches(2.2), Inches(12), Inches(1.2),
    "InvenPro", sz=54, bold=True, clr=WHITE, al=PP_ALIGN.CENTER)
txt(s, Inches(2), Inches(3.5), Inches(12), Inches(0.6),
    "Diagrama de Clases UML", sz=28, clr=ACCENT, al=PP_ALIGN.CENTER)
txt(s, Inches(2), Inches(4.4), Inches(12), Inches(0.5),
    "Sistema de Inventario y Facturacion para Minimercados", sz=14, clr=GRAY, al=PP_ALIGN.CENTER)
txt(s, Inches(2), Inches(6.8), Inches(12), Inches(0.5),
    "6 Modelos  ·  8 Servicios  ·  1 DataStore  ·  14 Componentes React",
    sz=12, clr=GRAY, al=PP_ALIGN.CENTER, fn=MONO)


# ═══════════════════════════════════════════════════
#  SLIDE 2 — Modelos 1/2 con relaciones entre sí
# ═══════════════════════════════════════════════════
s = nslide()
stitle(s, "Clases Modelo (1/2): Producto · Usuario · Cajero",
       "Entidades del dominio — constructor(data) + getters computados")
W = Inches(4.8)

cbox(s, Inches(0.4), Inches(1.2), W, "Producto", "entity", [
    "- sku: string", "- nombre: string", "- categoria: string",
    "- precio: number", "- costo: number", "- stock: number",
    "- min: number", "- vence: string | null",
    "- unidad: string", "- codigoBarras: string | null",
], [
    "+ constructor(data: object)",
    "+ get initials(): string",
    "+ get stockBajo(): boolean  // stock < min",
    "+ get valorTotal(): number  // stock * costo",
], hf=GREEN, brd=GREEN)

cbox(s, Inches(5.6), Inches(1.2), W, "Usuario", "entity", [
    "- usuario: string", "- nombre: string",
    "- rol: string       // Admin | Supervisor | Cajero",
    "- permisos: string[]", "- pass: string       // SHA-256",
], [
    "+ constructor(data: object)",
    "+ get initials(): string",
    "+ get esAdmin(): boolean  // rol==='Administrador'",
], hf=GREEN, brd=GREEN)

cbox(s, Inches(10.8), Inches(1.2), W, "Cajero", "entity", [
    "- id: string", "- nombre: string", "- doc: string",
    "- rol: string", "- estado: string    // activo | inactivo",
    "- turnoActivo: boolean | null",
    "- ingreso: string", "- ventas30d: number",
], [
    "+ constructor(data: object)",
    "+ get initials(): string",
    "+ get activo(): boolean  // estado==='activo'",
], hf=GREEN, brd=GREEN)

# Relación Cajero ─ ─ → Usuario (asociación: cajero tiene usuario vinculado)
arrow(s, Inches(10.8), Inches(3.0), Inches(10.4), Inches(3.0), clr=GREEN, w=1, dashed=True)
label(s, Inches(10.0), Inches(2.6), "tiene usuario\nasociado", clr=GREEN)


# ═══════════════════════════════════════════════════
#  SLIDE 3 — Modelos 2/2
# ═══════════════════════════════════════════════════
s = nslide()
stitle(s, "Clases Modelo (2/2): Proveedor · Turno · Factura",
       "Cada entidad se instancia desde filas camelizadas de Supabase")

cbox(s, Inches(0.4), Inches(1.2), W, "Proveedor", "entity", [
    "- id: string", "- nombre: string", "- nit: string",
    "- contacto: string", "- tel: string", "- email: string",
    "- ciudad: string", "- categoria: string",
    "- terminos: string    // Contado | 15d | 30d",
    "- estado: string      // activo | inactivo",
    "- ingresos: number", "- ultimoIngreso: string | null",
], [
    "+ constructor(data: object)",
    "+ get activo(): boolean  // estado==='activo'",
], hf=GREEN, brd=GREEN)

cbox(s, Inches(5.6), Inches(1.2), W, "Turno", "entity", [
    "- id: string", "- cajero: string",
    "- fechaIni: string", "- fechaFin: string | null",
    "- baseIni: number     // saldo apertura",
    "- ventas: number", "- transacciones: number",
    "- estado: string      // abierto | cerrado",
], [
    "+ constructor(data: object)",
    "+ get abierto(): boolean  // estado==='abierto'",
], hf=GREEN, brd=GREEN)

_fl, _ft, _fr, _fb, _fc, _fh = cbox(s, Inches(10.8), Inches(1.2), W, "Factura", "entity", [
    "- id: string", "- fecha: string", "- hora: string",
    "- cajero: string",
    "- metodo: string   // Efectivo|Nequi|...",
    "- total: number",
    "- items: FacturaItem[]",
], [
    "+ constructor(data: object)",
    "+ get cantidadItems(): number",
], hf=GREEN, brd=GREEN)

# Sub-clase implícita FacturaItem
fi_l = Inches(11.2)
fi_t = _fb + Inches(0.5)
cbox(s, fi_l, fi_t, Inches(4), "FacturaItem", "value object", [
    "- sku: string", "- nombre: string",
    "- q: number       // cantidad",
    "- precio: number",
], [], hf=RGBColor(0x22, 0xC5, 0x5E), brd=RGBColor(0x22, 0xC5, 0x5E), fs=7.5)

# Composición Factura ◆── FacturaItem
diamond_arrow(s, _fc, _fb, fi_l + Inches(2), fi_t, clr=GREEN, w=1.2)
label(s, Inches(12.0), _fb + Inches(0.12), "1..*", clr=GREEN)


# ═══════════════════════════════════════════════════
#  SLIDE 4 — Servicios 1/2
# ═══════════════════════════════════════════════════
s = nslide()
stitle(s, "Clases Servicio (1/2): Auth · Producto · Factura · Turno",
       "SRP — cada servicio gestiona una entidad. Stateless, operan contra Supabase.")
SW = Inches(3.65)

cbox(s, Inches(0.4), Inches(1.2), SW, "AuthService", "service", [
    "(stateless)",
], [
    "+ async login(usuario, pass)",
    "    → Usuario | null",
    "    // hash SHA-256, fallback migra",
    "+ async updatePassword(usr, newPass)",
    "    → Error | null",
], hf=PURPLE, brd=PURPLE)

cbox(s, Inches(4.35), Inches(1.2), SW, "ProductoService", "service", [
    "(stateless)",
], [
    "+ async getAll(): Producto[]",
    "+ async create(p): Error|null",
    "+ async update(sku, upd): Error|null",
    "+ async updateBarcode(sku, code)",
    "    → Error|null",
    "+ async incrementStock(sku, qty)",
    "    → Error|null  // RPC",
    "+ generateSku(): string  // P-XXXXX",
], hf=PURPLE, brd=PURPLE)

cbox(s, Inches(8.3), Inches(1.2), SW, "FacturaService", "service", [
    "(stateless)",
], [
    "+ async create(factura, items[])",
    "    → void",
    "    // 1. INSERT factura header",
    "    // 2. INSERT factura_items[]",
    "    // 3. RPC decrement_stock x item",
], hf=PURPLE, brd=PURPLE)

cbox(s, Inches(12.25), Inches(1.2), SW, "TurnoService", "service", [
    "(stateless)",
], [
    "+ async create(turno): void",
    "    // INSERT turnos (snakify)",
    "+ async close(id, updates): void",
    "    // UPDATE estado, ventas...",
], hf=PURPLE, brd=PURPLE)


# ═══════════════════════════════════════════════════
#  SLIDE 5 — Servicios 2/2
# ═══════════════════════════════════════════════════
s = nslide()
stitle(s, "Clases Servicio (2/2): Cajero · Proveedor · Ingreso · Config",
       "Operaciones CRUD asincronas — retornan Error | null para manejo de errores")

cbox(s, Inches(0.4), Inches(1.2), SW, "CajeroService", "service", [
    "(stateless)",
], [
    "+ async update(id, updates)",
    "    → Error|null  // cajeros",
    "+ async updateUsuario(usr, upd)",
    "    → Error|null  // usuarios_sistema",
], hf=PURPLE, brd=PURPLE)

cbox(s, Inches(4.35), Inches(1.2), SW, "ProveedorService", "service", [
    "(stateless)",
], [
    "+ async create(p): void",
    "    // INSERT proveedores (snakify)",
    "+ async update(id, updates): void",
    "    // UPDATE proveedores",
], hf=PURPLE, brd=PURPLE)

cbox(s, Inches(8.3), Inches(1.2), SW, "IngresoService", "service", [
    "(stateless)",
], [
    "+ async create(ingreso, detalle[])",
    "    → void",
    "    // INSERT header + detalle[]",
    "+ async update(id, hdr, det[])",
    "    → Error|null",
    "    // UPDATE hdr, DELETE+INSERT det",
    "    // recalcula totales",
], hf=PURPLE, brd=PURPLE)

cbox(s, Inches(12.25), Inches(1.2), SW, "ConfigService", "service", [
    "(stateless)",
], [
    "+ async save(clave, valor): void",
    "    // UPSERT configuracion",
    "    // actualiza MOCK cache",
    "+ async saveBatch(entries): void",
    "    // UPSERT multiples claves",
], hf=PURPLE, brd=PURPLE)


# ═══════════════════════════════════════════════════
#  SLIDE 6 — DataStore ◆── Modelos (COMPOSICIÓN)
# ═══════════════════════════════════════════════════
s = nslide()
stitle(s, "Relacion de Composicion: DataStore contiene Modelos",
       "\u25C6 Composicion (contiene) — DataStore.hydrate() crea instancias de cada modelo desde Supabase")

# Leyenda
rect(s, Inches(10), Inches(0.2), Inches(5.8), Inches(0.7), fill=WHITE, brd=BORDER)
txt(s, Inches(10.2), Inches(0.25), Inches(2), Pt(12), "Leyenda:", sz=8, bold=True, clr=DARK)
diamond_arrow(s, Inches(10.2), Inches(0.58), Inches(11.5), Inches(0.58), clr=TEAL, w=1.5)
txt(s, Inches(11.6), Inches(0.45), Inches(2), Pt(14), "Composicion (contiene 1..*)", sz=7.5, clr=DARK, fn=MONO)

# DataStore central
ds_x = Inches(5.8); ds_y = Inches(1.2); ds_w = Inches(4.8)
_dl, _dt, _dr, _db, _dc, _dh = cbox(s, ds_x, ds_y, ds_w, "DataStore", "singleton", [
    "- today: Date",
    "- productos: Producto[]",
    "- cajeros: Cajero[]",
    "- usuarios_sistema: Usuario[]",
    "- proveedores: Proveedor[]",
    "- turnos: Turno[]",
    "- facturas: Factura[]",
    "- ingresos: object[]",
    "- ventasMes: object[]",
    "- ventasCajero: object[]",
    "- topProductos: object[]",
    "- ventasHoy: object[]",
    "- configuracion: {clave: valor}",
], [
    "+ constructor()",
    "+ async hydrate(): void",
    "    // Promise.all(12 queries)",
    "    // camelize → new Modelo(row)",
], hf=TEAL, brd=TEAL)

# 6 Modelo boxes alrededor con composición arrows
modelos_izq = [
    ("Producto", ["- sku, nombre, categoria", "- precio, costo, stock, min", "- vence, unidad, codigoBarras"],
     ["+ get initials()", "+ get stockBajo()", "+ get valorTotal()"]),
    ("Usuario", ["- usuario, nombre, rol", "- permisos[], pass"],
     ["+ get initials()", "+ get esAdmin()"]),
    ("Cajero", ["- id, nombre, doc, rol", "- estado, turnoActivo", "- ingreso, ventas30d"],
     ["+ get initials()", "+ get activo()"]),
]
modelos_der = [
    ("Proveedor", ["- id, nombre, nit", "- contacto, tel, email", "- ciudad, cat, terminos", "- estado, ingresos"],
     ["+ get activo()"]),
    ("Turno", ["- id, cajero", "- fechaIni, fechaFin", "- baseIni, ventas, trans", "- estado"],
     ["+ get abierto()"]),
    ("Factura", ["- id, fecha, hora, cajero", "- metodo, total", "- items: FacturaItem[]"],
     ["+ get cantidadItems()"]),
]

MW = Inches(3.5)
# Izquierda
for i, (nm, att, met) in enumerate(modelos_izq):
    my = Inches(1.3) + Inches(i * 2.4)
    _ml, _mt, _mr, _mb, _mc, _mh = cbox(s, Inches(0.3), my, MW, nm, "entity", att, met, hf=GREEN, brd=GREEN, fs=7.5)
    # Flecha composición: DataStore ◆──→ Modelo (de derecha del modelo a izquierda del DataStore)
    arrow_y = my + _mh // 2
    diamond_arrow(s, ds_x, arrow_y, _mr + Pt(4), arrow_y, clr=TEAL, w=1.3)
    label(s, _mr + Pt(8), arrow_y - Pt(14), "1..*", sz=7, clr=TEAL)

# Derecha
for i, (nm, att, met) in enumerate(modelos_der):
    my = Inches(1.3) + Inches(i * 2.4)
    _ml, _mt, _mr, _mb, _mc, _mh = cbox(s, Inches(12.2), my, MW, nm, "entity", att, met, hf=GREEN, brd=GREEN, fs=7.5)
    arrow_y = my + _mh // 2
    diamond_arrow(s, _dr, arrow_y, _ml - Pt(4), arrow_y, clr=TEAL, w=1.3)
    label(s, _dr + Pt(8), arrow_y - Pt(14), "1..*", sz=7, clr=TEAL)

# Supabase abajo
sb_y = Inches(8.0)
rect(s, Inches(4), sb_y, Inches(8), Inches(0.65), fill=ORANGE_SOFT, brd=ORANGE)
txt(s, Inches(4), sb_y + Pt(3), Inches(8), Pt(14),
    "Supabase (PostgreSQL) — 12 tablas", sz=11, bold=True, clr=ORANGE, al=PP_ALIGN.CENTER)
txt(s, Inches(4), sb_y + Pt(20), Inches(8), Pt(12),
    "productos · cajeros · usuarios_sistema · proveedores · turnos · facturas · factura_items · ingresos · ingreso_detalle · configuracion",
    sz=6.5, clr=DARK, al=PP_ALIGN.CENTER, fn=MONO)
arrow(s, _dc, _db + Pt(4), _dc, sb_y, clr=ORANGE, w=1.5)
label(s, _dc + Pt(6), _db + Pt(12), "hydrate()\n12 queries", sz=7, clr=ORANGE)


# ═══════════════════════════════════════════════════
#  SLIDE 7 — Servicios ──→ Modelos (DEPENDENCIA)
# ═══════════════════════════════════════════════════
s = nslide()
stitle(s, "Relacion de Dependencia: Servicio opera sobre Modelo",
       "- - - → Dependencia — cada servicio realiza CRUD sobre su(s) modelo(s) via Supabase")

# Leyenda
rect(s, Inches(10.5), Inches(0.18), Inches(5.3), Inches(0.7), fill=WHITE, brd=BORDER)
txt(s, Inches(10.7), Inches(0.23), Inches(2), Pt(12), "Leyenda:", sz=8, bold=True, clr=DARK)
arrow(s, Inches(10.7), Inches(0.56), Inches(12), Inches(0.56), clr=PURPLE, w=1.2, dashed=True)
txt(s, Inches(12.1), Inches(0.42), Inches(3), Pt(14), "Dependencia (usa / opera sobre)", sz=7.5, clr=DARK, fn=MONO)

# Servicios a la izquierda (compactos)
svcs = [
    ("AuthService",      ["+ login(usr, pass): Usuario?", "+ updatePassword(usr, p)"]),
    ("ProductoService",  ["+ getAll(): Producto[]", "+ create(p), update(sku, u)", "+ updateBarcode(), incrementStock()", "+ generateSku(): string"]),
    ("FacturaService",   ["+ create(factura, items[])", "  // decrement_stock RPC"]),
    ("TurnoService",     ["+ create(turno)", "+ close(id, updates)"]),
    ("CajeroService",    ["+ update(id, upd)", "+ updateUsuario(usr, upd)"]),
    ("ProveedorService", ["+ create(p)", "+ update(id, upd)"]),
    ("IngresoService",   ["+ create(ingreso, det[])", "+ update(id, hdr, det[])"]),
    ("ConfigService",    ["+ save(clave, valor)", "+ saveBatch(entries)"]),
]

# Modelos a la derecha
mdls = [
    ("Usuario",    ["- usuario, nombre", "- rol, permisos, pass"], ["+ get esAdmin()"]),
    ("Producto",   ["- sku, nombre, precio", "- stock, min, vence"], ["+ get stockBajo()"]),
    ("Factura",    ["- id, fecha, cajero", "- metodo, total, items[]"], ["+ get cantidadItems()"]),
    ("Turno",      ["- id, cajero, estado", "- fechaIni, ventas"], ["+ get abierto()"]),
    ("Cajero",     ["- id, nombre, doc", "- estado, ventas30d"], ["+ get activo()"]),
    ("Proveedor",  ["- id, nombre, nit", "- estado, terminos"], ["+ get activo()"]),
]

svc_x = Inches(0.3); svc_w = Inches(3.8)
mdl_x = Inches(11.8); mdl_w = Inches(3.5)

svc_centers = {}
mdl_centers = {}

for i, (nm, meths) in enumerate(svcs):
    sy = Inches(1.15) + Inches(i * 0.93)
    _l, _t, _r, _b, _c, _h = cbox(s, svc_x, sy, svc_w, nm, None, [], meths, hf=PURPLE, brd=PURPLE, fs=7)
    svc_centers[nm] = (_r, sy + _h // 2)

for i, (nm, attrs, meths) in enumerate(mdls):
    my = Inches(1.15) + Inches(i * 1.2)
    _l, _t, _r, _b, _c, _h = cbox(s, mdl_x, my, mdl_w, nm, None, attrs, meths, hf=GREEN, brd=GREEN, fs=7)
    mdl_centers[nm] = (_l, my + _h // 2)

# Flechas de dependencia servicio → modelo
deps = [
    ("AuthService", "Usuario"),
    ("ProductoService", "Producto"),
    ("FacturaService", "Factura"),
    ("TurnoService", "Turno"),
    ("CajeroService", "Cajero"),
    ("ProveedorService", "Proveedor"),
]
for svc_nm, mdl_nm in deps:
    sx, sy = svc_centers[svc_nm]
    mx, my = mdl_centers[mdl_nm]
    arrow(s, sx + Pt(4), sy, mx - Pt(4), my, clr=PURPLE, w=1, dashed=True)

# Flechas cruzadas
# CajeroService → Usuario
sx, sy = svc_centers["CajeroService"]
mx, my = mdl_centers["Usuario"]
arrow(s, sx + Pt(4), sy - Pt(8), mx - Pt(4), my + Pt(8), clr=PURPLE, w=0.8, dashed=True)

# IngresoService → Producto
sx, sy = svc_centers["IngresoService"]
mx, my = mdl_centers["Producto"]
arrow(s, sx + Pt(4), sy, mx - Pt(4), my + Pt(8), clr=PURPLE, w=0.8, dashed=True)

# Etiquetas centrales
txt(s, Inches(5), Inches(1.0), Inches(6), Pt(14),
    "- - - - - - - - - -  opera sobre  - - - - - - - - - -",
    sz=9, clr=PURPLE, al=PP_ALIGN.CENTER, fn=MONO)

# ConfigService → configuracion (sin modelo)
sx, sy = svc_centers["ConfigService"]
rect(s, Inches(11.8), Inches(8.0), mdl_w, Inches(0.45), fill=ORANGE_SOFT, brd=ORANGE)
txt(s, Inches(11.8), Inches(8.05), mdl_w, Pt(14), "configuracion{}", sz=9, bold=True, clr=ORANGE, al=PP_ALIGN.CENTER, fn=MONO)
arrow(s, sx + Pt(4), sy, Inches(11.8) - Pt(4), Inches(8.2), clr=PURPLE, w=0.8, dashed=True)


# ═══════════════════════════════════════════════════
#  SLIDE 8 — Componentes ──→ Servicios (USO)
# ═══════════════════════════════════════════════════
s = nslide()
stitle(s, "Relacion de Uso: Componente consume Servicio",
       "──→ Uso (dependency) — cada componente llama metodos de DB.servicio")

rect(s, Inches(10.5), Inches(0.18), Inches(5.3), Inches(0.7), fill=WHITE, brd=BORDER)
txt(s, Inches(10.7), Inches(0.23), Inches(2), Pt(12), "Leyenda:", sz=8, bold=True, clr=DARK)
arrow(s, Inches(10.7), Inches(0.56), Inches(12), Inches(0.56), clr=ACCENT, w=1.2)
txt(s, Inches(12.1), Inches(0.42), Inches(3), Pt(14), "Uso (llama metodos de)", sz=7.5, clr=DARK, fn=MONO)

# Componentes a la izquierda
comps = [
    ("Login",       ["- user, pass, error"],          ["+ submit() → DB.auth.login()"]),
    ("ShiftOpen",   ["- base, caja, loading"],         ["+ handleOpen() → DB.turnos.create()"]),
    ("POS",         ["- cart[], total, method"],        ["+ finalizePayment()", "  → DB.facturas.create()", "  → DB.turnos.close()"]),
    ("Inventario",  ["- productos[], editing"],         ["+ guardarProducto()", "  → DB.productos.update()", "+ asignarCodigo()", "  → DB.productos.updateBarcode()"]),
    ("Ingreso",     ["- items[], origen, prov"],        ["+ guardarIngreso()", "  → DB.productos.create()", "  → DB.productos.incrementStock()", "  → DB.ingresos.create()"]),
    ("Proveedores", ["- list[], editing"],              ["+ guardar()", "  → DB.proveedores.create()/update()"]),
    ("Cajeros",     ["- tab, cfgCajero"],               ["+ handleSave()", "  → DB.cajeros.update()", "  → DB.cajeros.updateUsuario()"]),
    ("Ajustes",     ["- apiKey, urlApi"],                ["+ guardar() → DB.config.saveBatch()", "+ borrarKey() → DB.config.save()"]),
]

# Servicios a la derecha
svc_list = [
    ("AuthService",      ["+ login()", "+ updatePassword()"]),
    ("TurnoService",     ["+ create()", "+ close()"]),
    ("FacturaService",   ["+ create(factura, items)"]),
    ("ProductoService",  ["+ create(), update()", "+ updateBarcode()", "+ incrementStock()", "+ generateSku()"]),
    ("IngresoService",   ["+ create()", "+ update()"]),
    ("ProveedorService", ["+ create()", "+ update()"]),
    ("CajeroService",    ["+ update()", "+ updateUsuario()"]),
    ("ConfigService",    ["+ save()", "+ saveBatch()"]),
]

comp_x = Inches(0.3); comp_w = Inches(4.2)
sv_x = Inches(11.5); sv_w = Inches(3.8)

comp_pos = {}
svc_pos = {}

for i, (nm, attrs, meths) in enumerate(comps):
    cy = Inches(1.1) + Inches(i * 0.95)
    _l, _t, _r, _b, _c, _h = cbox(s, comp_x, cy, comp_w, nm, None, attrs, meths, hf=ACCENT, brd=ACCENT, fs=7)
    comp_pos[nm] = (_r, cy + _h // 2)

for i, (nm, meths) in enumerate(svc_list):
    sy = Inches(1.1) + Inches(i * 0.95)
    _l, _t, _r, _b, _c, _h = cbox(s, sv_x, sy, sv_w, nm, None, [], meths, hf=PURPLE, brd=PURPLE, fs=7)
    svc_pos[nm] = (_l, sy + _h // 2)

# Flechas de uso
use_map = [
    ("Login",       ["AuthService"]),
    ("ShiftOpen",   ["TurnoService"]),
    ("POS",         ["FacturaService", "TurnoService"]),
    ("Inventario",  ["ProductoService"]),
    ("Ingreso",     ["ProductoService", "IngresoService"]),
    ("Proveedores", ["ProveedorService"]),
    ("Cajeros",     ["CajeroService"]),
    ("Ajustes",     ["ConfigService"]),
]

for comp_nm, svc_nms in use_map:
    cx, cy = comp_pos[comp_nm]
    for svc_nm in svc_nms:
        sx, sy = svc_pos[svc_nm]
        arrow(s, cx + Pt(4), cy, sx - Pt(4), sy, clr=ACCENT, w=1)

txt(s, Inches(5.5), Inches(1.0), Inches(5), Pt(14),
    "- - - - - - -  usa  - - - - - - -",
    sz=9, clr=ACCENT, al=PP_ALIGN.CENTER, fn=MONO)


# ═══════════════════════════════════════════════════
#  SLIDE 9 — DataStore + Globals
# ═══════════════════════════════════════════════════
s = nslide()
stitle(s, "DataStore: Almacen Central + Instancias Globales",
       "window.MOCK = DataStore (datos)  ·  window.DB = { servicios } (operaciones)")

ds_w2 = Inches(6)
cbox(s, Inches(0.4), Inches(1.2), ds_w2, "DataStore", "singleton", [
    "- today: Date",
    "- productos: Producto[]",
    "- cajeros: Cajero[]",
    "- usuarios_sistema: Usuario[]",
    "- proveedores: Proveedor[]",
    "- turnos: Turno[]",
    "- facturas: Factura[]       // con items[] anidados",
    "- ingresos: object[]        // con detalle[] anidado",
    "- ventasMes: object[]",
    "- ventasCajero: object[]",
    "- topProductos: object[]",
    "- ventasHoy: object[]",
    "- configuracion: {clave: valor}",
], [
    "+ constructor()   // inicializa arrays vacios",
    "+ async hydrate(): void",
    "    // 1. Promise.all(12 queries Supabase)",
    "    // 2. camelize(snake_case → camelCase)",
    "    // 3. new Producto(row) por cada fila",
    "    // 4. new Cajero(row), new Usuario(row)...",
    "    // 5. Factura: anida items[], ordena fecha",
    "    // 6. Ingreso: anida detalle[]",
    "    // 7. configuracion filas → mapa {k:v}",
], hf=TEAL, brd=TEAL, fs=8)

# window.DB box
db_x = Inches(7.5); db_y = Inches(1.2); db_w = Inches(4.5)
rect(s, db_x, db_y, db_w, Inches(0.4), fill=PURPLE, brd=PURPLE)
txt(s, db_x, db_y + Pt(4), db_w, Pt(18), "window.DB", sz=12, bold=True, clr=WHITE, al=PP_ALIGN.CENTER)

svc_items = [
    "auth: AuthService",
    "productos: ProductoService",
    "facturas: FacturaService",
    "turnos: TurnoService",
    "cajeros: CajeroService",
    "proveedores: ProveedorService",
    "ingresos: IngresoService",
    "config: ConfigService",
]
for i, item in enumerate(svc_items):
    iy = db_y + Inches(0.5) + Inches(i * 0.42)
    rect(s, db_x, iy, db_w, Inches(0.35), fill=PURPLE_SOFT, brd=PURPLE)
    txt(s, db_x + Pt(10), iy + Pt(4), db_w - Pt(20), Pt(14), item,
        sz=9, clr=PURPLE, fn=MONO, al=PP_ALIGN.LEFT)

# window.MOCK label
rect(s, Inches(0.4), Inches(7.6), Inches(3), Inches(0.4), fill=TEAL, brd=TEAL)
txt(s, Inches(0.4), Inches(7.62), Inches(3), Pt(18), "window.MOCK = DataStore",
    sz=10, bold=True, clr=WHITE, al=PP_ALIGN.CENTER)

# Supabase
rect(s, Inches(7.5), Inches(7.6), db_w, Inches(0.72), fill=ORANGE_SOFT, brd=ORANGE)
txt(s, Inches(7.5), Inches(7.63), db_w, Pt(14), "Supabase (PostgreSQL)",
    sz=10, bold=True, clr=ORANGE, al=PP_ALIGN.CENTER)
txt(s, Inches(7.5), Inches(7.95), db_w, Pt(12),
    "12 tablas + RPCs", sz=8, clr=DARK, al=PP_ALIGN.CENTER, fn=MONO)

# Arrows
arrow(s, Inches(3.4), Inches(7.2), Inches(3.4), Inches(7.6), clr=TEAL, w=1.5)
arrow(s, Inches(9.75), Inches(7.2), Inches(9.75), Inches(7.6), clr=ORANGE, w=1.5)
label(s, Inches(9.9), Inches(7.1), "CRUD async", sz=7, clr=ORANGE)

# Utilidades globales
rect(s, Inches(13), Inches(1.2), Inches(2.6), Inches(3.5), fill=TEAL_SOFT, brd=TEAL)
txt(s, Inches(13.1), Inches(1.3), Inches(2.4), Pt(14), "Helpers Globales", sz=9, bold=True, clr=TEAL)
mtxt(s, Inches(13.1), Inches(1.65), Inches(2.4), [
    "camelize(obj)",
    "snakify(obj)",
    "hashPass(plain)",
    "fmtCOP(n)",
    "daysFromNow(str)",
    "exportXlsx(f, sh)",
    "hydrateData()",
], sz=7.5, clr=DARK)


# ═══════════════════════════════════════════════════
#  SLIDE 10 — Componentes React con relaciones
# ═══════════════════════════════════════════════════
s = nslide()
stitle(s, "Componentes React: App → Flujo Cajero + Panel Admin",
       "Componentes funcionales con hooks — App renderiza condicionalmente segun stage y rol")

# App root
app_w = Inches(6)
_al, _at, _ar, _ab, _ac, _ah = cbox(s, Inches(5), Inches(1.1), app_w, "App", "root", [
    "- theme: 'light'|'dark'",
    "- stage: login|shift-open|pos|shift-closed|admin",
    "- user: {nombre, rol, usuario, permisos}",
    "- shift: Turno + metadata", "- adminPage: string",
], [
    "+ onLogin(u) — hydrateData(), set stage",
    "+ onLogout() — clear localStorage",
], hf=RGBColor(0x1D, 0x4E, 0xD8), brd=RGBColor(0x1D, 0x4E, 0xD8))

# Cashier flow (left)
txt(s, Inches(0.3), Inches(3.7), Inches(3), Pt(16), "Flujo Cajero", sz=11, bold=True, clr=ACCENT)
c_comps = [
    ("Login",      ["- user, pass, error"],    ["+ submit() → DB.auth.login()"]),
    ("ShiftOpen",  ["- base, caja"],            ["+ handleOpen() → DB.turnos.create()"]),
    ("POS",        ["- cart[], total, method"],  ["+ finalizePayment()", "  → DB.facturas.create()"]),
    ("ShiftClosed",["- summary"],               ["+ render(resumen turno)"]),
]
ccw = Inches(3.8)
prev_bottom = None
for i, (nm, att, met) in enumerate(c_comps):
    cy = Inches(4.0) + Inches(i * 1.2)
    _l, _t, _r, _b, _c, _h = cbox(s, Inches(0.3), cy, ccw, nm, None, att, met, hf=ACCENT, brd=ACCENT, fs=7)
    if prev_bottom:
        arrow(s, Inches(0.3) + ccw//2, prev_bottom + Pt(2), Inches(0.3) + ccw//2, cy - Pt(2), clr=ACCENT, w=1)
    prev_bottom = _b

# Arrow from App to Login
arrow(s, _ac - Inches(1.5), _ab + Pt(2), Inches(0.3) + ccw//2, Inches(4.0) - Pt(2), clr=ACCENT, w=1.2)
label(s, Inches(2.5), Inches(3.45), "renders", clr=ACCENT)

# Admin flow (right)
txt(s, Inches(5), Inches(3.7), Inches(5), Pt(16), "Panel Admin", sz=11, bold=True, clr=ACCENT)
a_comps = [
    ("Dashboard",   "MOCK (lectura)"),
    ("Inventario",  "ProductoService"),
    ("Ingreso",     "IngresoService, ProductoService"),
    ("Proveedores", "ProveedorService"),
    ("Cajeros",     "CajeroService"),
    ("Reportes",    "MOCK (lectura)"),
    ("Ajustes",     "ConfigService"),
]
acw = Inches(2.6)
agap = Inches(0.18)

for i, (nm, svc) in enumerate(a_comps):
    col = i % 4
    row = i // 4
    ax = Inches(5) + col * (acw + agap)
    ay = Inches(4.0) + row * Inches(2.2)
    rect(s, ax, ay, acw, Inches(0.7), fill=ACCENT_SOFT, brd=ACCENT)
    txt(s, ax, ay + Pt(4), acw, Pt(14), nm, sz=10, bold=True, clr=ACCENT, al=PP_ALIGN.CENTER)
    txt(s, ax, ay + Pt(22), acw, Pt(20), svc, sz=7, clr=PURPLE, al=PP_ALIGN.CENTER, fn=MONO)

# Arrow from App to admin
arrow(s, _ac + Inches(1.5), _ab + Pt(2), Inches(8), Inches(4.0) - Pt(2), clr=ACCENT, w=1.2)
label(s, Inches(8.2), Inches(3.45), "renders", clr=ACCENT)

# Sidebar
rect(s, Inches(5), Inches(6.9), Inches(2.6), Inches(0.55), fill=GREEN_SOFT, brd=GREEN)
txt(s, Inches(5), Inches(6.93), Inches(2.6), Pt(14), "Sidebar", sz=10, bold=True, clr=GREEN, al=PP_ALIGN.CENTER)
txt(s, Inches(5), Inches(7.15), Inches(2.6), Pt(12), "navegacion + logout", sz=7, clr=DARK, al=PP_ALIGN.CENTER, fn=MONO)

# UI compartida
ui_names = ["Icon", "Modal", "Toast", "Pagination", "Spark"]
for i, nm in enumerate(ui_names):
    ux = Inches(8) + Inches(i * 1.1)
    rect(s, ux, Inches(7.5), Inches(1), Inches(0.4), fill=GREEN_SOFT, brd=GREEN)
    txt(s, ux, Inches(7.52), Inches(1), Pt(14), nm, sz=7.5, bold=True, clr=GREEN, al=PP_ALIGN.CENTER)
txt(s, Inches(8), Inches(7.15), Inches(5.5), Pt(14), "UI Compartida (usados por todos)",
    sz=9, bold=True, clr=GREEN)


# ═══════════════════════════════════════════════════
#  SLIDE 11 — Diagrama completo de relaciones
# ═══════════════════════════════════════════════════
s = nslide(DARK)
txt(s, Inches(0.5), Inches(0.2), Inches(10), Inches(0.45),
    "Diagrama Completo de Relaciones", sz=21, bold=True, clr=WHITE)

# Legend
for i, (lb, c, bg) in enumerate([
    ("Modelo", GREEN, GREEN_SOFT), ("Servicio", PURPLE, PURPLE_SOFT),
    ("DataStore", TEAL, TEAL_SOFT), ("Componente", ACCENT, ACCENT_SOFT),
    ("Supabase", ORANGE, ORANGE_SOFT),
]):
    lx = Inches(0.5) + Inches(i * 2.8)
    rect(s, lx, Inches(0.72), Inches(2.5), Inches(0.28), fill=bg, brd=c)
    txt(s, lx, Inches(0.73), Inches(2.5), Pt(14), lb, sz=8.5, bold=True, clr=c, al=PP_ALIGN.CENTER)

uml = (
    "                                          ┌───────────────────────────────────────────────┐\n"
    "                                          │              App (React Root)                 │\n"
    "                                          │   stage · user · shift · theme · adminPage    │\n"
    "                                          └──────────┬──────────────────┬─────────────────┘\n"
    "                                 ┌───────────────────┘   renders        └─────────────────────────┐\n"
    "                      ┌──────────┴──────────┐                              ┌─────────────────────┴────────┐\n"
    "                      │    Flujo Cajero     │                              │        Panel Admin           │\n"
    "                      │ Login → ShiftOpen   │                              │  Dashboard · Inventario      │\n"
    "                      │ → POS → ShiftClosed │                              │  Ingreso · Proveedores       │\n"
    "                      └──────────┬──────────┘                              │  Cajeros · Reportes · Ajustes│\n"
    "                                 │                                         └──────────────┬───────────────┘\n"
    "                                 │                  usa                                   │\n"
    "                                 └──────────────────────┐        ┌────────────────────────┘\n"
    "                                                        ▼        ▼\n"
    "                                 ┌─────────────────────────────────────────────────────────┐\n"
    "                                 │                    window.DB                            │\n"
    "                                 │  ┌──────────┐ ┌──────────────┐ ┌──────────────┐        │\n"
    "                                 │  │   auth   │ │  productos   │ │   facturas   │  ...   │\n"
    "                                 │  └──────────┘ └──────────────┘ └──────────────┘        │\n"
    "                                 │  ┌──────────┐ ┌──────────────┐ ┌──────────────┐        │\n"
    "                                 │  │  turnos  │ │  proveedores │ │   ingresos   │ config │\n"
    "                                 │  └──────────┘ └──────────────┘ └──────────────┘        │\n"
    "                                 └────────────────────────┬────────────────────────────────┘\n"
    "                                                          │ opera sobre\n"
    "                                                          ▼\n"
    "                                 ┌─────────────────────────────────────────────────────────┐\n"
    "                                 │                window.MOCK (DataStore)                  │\n"
    "                                 │ ◆─ Producto[]  ◆─ Usuario[]   ◆─ Cajero[]              │\n"
    "                                 │ ◆─ Proveedor[] ◆─ Turno[]    ◆─ Factura[]             │\n"
    "                                 │    ventasMes[] · ventasCajero[] · configuracion{}       │\n"
    "                                 └────────────────────────┬────────────────────────────────┘\n"
    "                                                          │ hydrate() — 12 queries\n"
    "                                                          ▼\n"
    "                                 ┌─────────────────────────────────────────────────────────┐\n"
    "                                 │               Supabase (PostgreSQL)                     │\n"
    "                                 │  productos · cajeros · usuarios_sistema · proveedores   │\n"
    "                                 │  turnos · facturas · factura_items · ingresos           │\n"
    "                                 │  ingreso_detalle · configuracion · RPCs                 │\n"
    "                                 └─────────────────────────────────────────────────────────┘"
)
txt(s, Inches(0.2), Inches(1.15), Inches(15.6), Inches(7.5), uml,
    sz=8, clr=RGBColor(0xA0, 0xAE, 0xC0), fn=MONO)


# ═══════════════════════════════════════════════════
#  SLIDE 12 — Tabla de dependencias
# ═══════════════════════════════════════════════════
s = nslide()
stitle(s, "Matriz de Dependencias: Componente → Servicio → Modelo")

rows_data = [
    ["Componente",  "Servicio(s)",                   "Modelo(s)",               "Metodos que invoca"],
    ["Login",       "AuthService",                   "Usuario",                 "DB.auth.login(usr, pass)"],
    ["ShiftOpen",   "TurnoService",                  "Turno",                   "DB.turnos.create(turnoRow)"],
    ["POS",         "FacturaService\nTurnoService",  "Factura, Turno",          "DB.facturas.create(f, items)\nDB.turnos.close(id, upd)"],
    ["ShiftClosed", "TurnoService",                  "Turno",                   "DB.turnos.close(id, upd)"],
    ["Dashboard",   "(lectura MOCK)",                "ventasMes, ventasHoy",    "(sin llamadas a DB)"],
    ["Inventario",  "ProductoService",               "Producto",                "DB.productos.update(sku, upd)\nDB.productos.updateBarcode(sku, c)"],
    ["Ingreso",     "IngresoService\nProductoService","Producto",               "DB.productos.create(p)\nDB.productos.incrementStock(s,q)\nDB.ingresos.create(ing, det)"],
    ["Proveedores", "ProveedorService",              "Proveedor",               "DB.proveedores.create(p)\nDB.proveedores.update(id, upd)"],
    ["Cajeros",     "CajeroService",                 "Cajero, Usuario",         "DB.cajeros.update(id, upd)\nDB.cajeros.updateUsuario(u, upd)"],
    ["Reportes",    "(lectura MOCK)",                "Factura",                 "(filtros + exportXlsx)"],
    ["Ajustes",     "ConfigService",                 "configuracion{}",         "DB.config.saveBatch(entries)\nDB.config.save(clave, valor)"],
]

tbl = s.shapes.add_table(len(rows_data), 4, Inches(0.5), Inches(1.0), Inches(15), Inches(7.5)).table
tbl.columns[0].width = Inches(2)
tbl.columns[1].width = Inches(3.3)
tbl.columns[2].width = Inches(3.2)
tbl.columns[3].width = Inches(6.5)

for r, row in enumerate(rows_data):
    for c, val in enumerate(row):
        cell = tbl.cell(r, c)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(9.5 if r > 0 else 10.5)
            p.font.name = MONO if (r > 0 and c > 0) else FONT
            p.font.bold = (r == 0)
            p.font.color.rgb = WHITE if r == 0 else DARK
        if r == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT
        elif r % 2 == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = STRIPE


# ═══════════════════════════════════════
output = r"C:\invenpro\InvenPro_Diagrama_Clases_UML.pptx"
prs.save(output)
print(f"OK -> {output}")
print(f"Total: {len(prs.slides)} slides")
